from pptx import Presentation

def extract_and_format_specific_info(pptx_path, slide_index, table_index, target_label):
    # Charger la présentation PowerPoint
    presentation = Presentation(pptx_path)

    # Sélectionner la diapositive (index de la diapositive)
    slide = presentation.slides[slide_index - 1]  # Indice de la diapositive (commence à 0)

    try:
        # Sélectionner le tableau sur la diapositive (index du tableau)
        shape = slide.shapes[table_index]

        # Vérifier si le Shape à l'index spécifié est un tableau
        if shape.has_table:
            table = shape.table

            # Parcourir les lignes du tableau
            for row in table.rows:
                # Ignorer les lignes sans cellules
                if not row.cells:
                    continue

                # Vérifier si la première cellule contient le libellé cible
                if target_label.lower() in row.cells[0].text.lower():
                    # Construire le résultat avec le libellé formaté
                    formatted_info = f"{target_label} : {row.cells[1].text.strip()}"
                    
                    # Afficher le résultat
                    return formatted_info

            # Si le libellé cible n'est pas trouvé
            print(f"Le libellé cible '{target_label}' n'a pas été trouvé dans le tableau.")
            return None

        else:
            print(f"L'objet à l'index {table_index} sur la diapositive {slide_index} n'est pas un tableau.")
            return None

    except IndexError:
        print(f"Aucun objet trouvé à l'index {table_index} sur la diapositive {slide_index}.")
        return None

# Exemple d'utilisation pour la quatrième diapositive (slide 4), tableau d'index 1
pptx_file_path = 'Exemple FASEP.pptx'
slide_index_to_extract = 4  # Indice de la diapositive que vous souhaitez extraire (commence à 1)
table_index_to_extract = 1  # Indice du tableau que vous souhaitez extraire
labels_to_extract = "Avis sur le versement intermédiaire"

avis = extract_and_format_specific_info(pptx_file_path, slide_index_to_extract, table_index_to_extract, labels_to_extract)

