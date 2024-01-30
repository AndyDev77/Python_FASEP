# main.py
from pptx import Presentation

# Configuration
pptx_file_path = 'Exemple FASEP.pptx'

# Fonction pour extraire les dates du tableau
def extract_dates_from_table(slide_index, table_index, labels_to_extract):
    presentation = Presentation(pptx_file_path)
    slide = presentation.slides[slide_index - 1]
    table = slide.shapes[table_index].table
    dates = []

    for row in table.rows:
        if not row.cells:
            continue
        if any(label.lower() in row.cells[0].text.lower() for label in labels_to_extract):
            date_str = row.cells[1].text.strip()
            dates.append(date_str)

    return dates

# Fonction pour extraire le montant du tableau
def extract_montant_from_table(slide_index, table_index, target_label):
    presentation = Presentation(pptx_file_path)
    slide = presentation.slides[slide_index - 1]

    try:
        shape = slide.shapes[table_index]

        if shape.has_table:
            table = shape.table
            for row in table.rows:
                if not row.cells:
                    continue
                if target_label.lower() in row.cells[0].text.lower():
                    subvention_amount = f"{target_label}: {row.cells[1].text.strip()}"
                    return subvention_amount
            print(f"Le libellé cible '{target_label}' n'a pas été trouvé dans le tableau.")
            return None

        else:
            print(f"L'objet à l'index {table_index} sur la diapositive {slide_index} n'est pas un tableau.")
            return None

    except IndexError:
        print(f"Aucun objet trouvé à l'index {table_index} sur la diapositive {slide_index}.")
        return None

# Fonction pour extraire l'avis du tableau
def extract_avis_from_table(slide_index, table_index, target_label):
    presentation = Presentation(pptx_file_path)
    slide = presentation.slides[slide_index - 1]

    try:
        shape = slide.shapes[table_index]

        if shape.has_table:
            table = shape.table
            for row in table.rows:
                if not row.cells:
                    continue
                if target_label.lower() in row.cells[0].text.lower():
                    formatted_info = f"{target_label} : {row.cells[1].text.strip()}"
                    return formatted_info
            print(f"Le libellé cible '{target_label}' n'a pas été trouvé dans le tableau.")
            return None

        else:
            print(f"L'objet à l'index {table_index} sur la diapositive {slide_index} n'est pas un tableau.")
            return None

    except IndexError:
        print(f"Aucun objet trouvé à l'index {table_index} sur la diapositive {slide_index}.")
        return None

# Fonction principale
def main():
    # Diapositive 3, Tableau 1
    slide_index_dates = 3
    table_index_dates = 1
    labels_dates = ["Date de signature", "Montant et date de paiement de l'acompte"]

    # Diapositive 1, Tableau 2
    slide_index_montant = 1
    table_index_montant = 2
    labels_montant = "Montant du FASEP"

    # Diapositive 4, Tableau 1
    slide_index_avis = 4
    table_index_avis = 1
    labels_avis = "Avis sur le versement intermédiaire"

    # Extraction des dates
    dates_result = extract_dates_from_table(slide_index_dates, table_index_dates, labels_dates)
    if dates_result:
        print(f"Dates de signature de la convention de don pour la subvention: {dates_result}")

    # Extraction du montant de la subvention
    montant_result = extract_montant_from_table(slide_index_montant, table_index_montant, labels_montant)
    if montant_result:
        print(f"{montant_result}")

    # Extraction de l'avis
    avis_result = extract_avis_from_table(slide_index_avis, table_index_avis, labels_avis)
    if avis_result:
        print(f"{avis_result}")

if __name__ == "__main__":
    main()
