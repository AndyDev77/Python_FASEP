from pptx import Presentation

def extract_dates_from_table(pptx_path, slide_index, table_index, labels_to_extract):
    # Charger la présentation PowerPoint
    presentation = Presentation(pptx_path)

    # Sélectionner la diapositive (index de la diapositive)
    slide = presentation.slides[slide_index - 1]  # Indice de la diapositive (commence à 0)

    # Sélectionner le tableau sur la diapositive (index du tableau)
    table = slide.shapes[table_index].table

    # Initialiser une liste pour stocker les dates
    dates = []

    for row in table.rows:
        if not row.cells:
            continue

        if any(label.lower() in row.cells[0].text.lower() for label in labels_to_extract):
            date_str = row.cells[1].text.strip()
            dates.append(date_str)

    return dates

# Exemple d'utilisation
pptx_file_path = 'Exemple FASEP.pptx'
slide_index_to_extract = 3  # Indice de la diapositive que vous souhaitez extraire (commence à 1)
table_index_to_extract = 1  # Indice du tableau que vous souhaitez extraire
labels_to_extract = ["Date de signature", "Montant et date de paiement de l'acompte"]

dates = extract_dates_from_table(pptx_file_path, slide_index_to_extract, table_index_to_extract, labels_to_extract)

# Afficher les dates
# print("Dates de signature de la convention de don pour la subvention:", dates)
